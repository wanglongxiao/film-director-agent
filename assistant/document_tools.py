"""文档生成/读取工具 —— 在 AgentKit sandbox 中真实生成并读回 word/pdf/ppt/html。

背景与设计取舍：
- 账号当前的 AgentKit sandbox（AGENTKIT_TOOL_ID*）是「纯代码环境（CodeEnv）」，
  沙箱内并没有 skill 中心的技能运行器（/home/gem/veadk_skills/agent.py），因此
  内置的 execute_skills（本质是在沙箱里执行 `python agent.py <prompt> --skills ...`）
  无法加载 skill 中心的 docx / pdf-processing-pro 技能包。仅设置 SKILL_SPACE_ID 也无效，
  因为缺的是「运行器 + 技能目录」这套 SkillEnv 基础设施（需要新建云端 SkillEnv 工具）。
- 但该沙箱本身可联网、可执行 Python，且镜像已内置 python-docx / python-pptx / pypdf /
  matplotlib 与 weasyprint CLI。因此这里用「沙箱兜底」方案：把 word/pdf/ppt/html 的
  「生成 + 读取」封装成两个可被 agent 直接调用的工具，全部在同一个 AgentKit sandbox
  内用 RunCode 真实执行，无需依赖 skill 中心，也无需临时 pip 安装。

对外暴露两个工具（会作为 FunctionTool 挂到 root_agent）：
- create_document(doc_format, filename, content, title=None, tool_context=None)
- read_document(path, doc_format=None, tool_context=None)
"""

from __future__ import annotations

import base64
import json
import os
from typing import Optional

from google.adk.tools import ToolContext

from veadk.tools.builtin_tools._agentkit import (
    invoke_agentkit_run_code,
    resolve_agentkit_tool_id,
)
from veadk.utils.logger import get_logger

logger = get_logger(__name__)

# 沙箱内文档产物统一存放目录（供后续 read_document 读回）。
_SANDBOX_DOC_DIR = "/home/gem/veadk_docs"
_SUPPORTED_FORMATS = ("docx", "pdf", "pptx", "html")
_RUN_TIMEOUT = 300


# --- 在沙箱里真正执行的脚本（自包含，依赖镜像内置库） --------------------------------
# 通过 base64 传入 params，彻底规避引号/转义问题。params 结构：
#   {"action":"create"|"read","fmt":...,"filename":...,"path":...,
#    "content":...,"title":...,"base_dir":...}
_SANDBOX_SCRIPT = r'''
import base64, json, os, subprocess

params = json.loads(base64.b64decode("__PARAMS_B64__").decode("utf-8"))


def _esc(s):
    return (s or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _lines_to_html_body(title, content):
    parts = []
    if title:
        parts.append("<h1>%s</h1>" % _esc(title))
    for raw in (content or "").split("\n"):
        line = raw.rstrip()
        if not line.strip():
            continue
        if line.startswith("## "):
            parts.append("<h2>%s</h2>" % _esc(line[3:].strip()))
        elif line.startswith("# "):
            parts.append("<h2>%s</h2>" % _esc(line[2:].strip()))
        elif line.startswith("- ") or line.startswith("* "):
            parts.append("<li>%s</li>" % _esc(line[2:].strip()))
        else:
            parts.append("<p>%s</p>" % _esc(line))
    html = []
    in_ul = False
    for p in parts:
        if p.startswith("<li>"):
            if not in_ul:
                html.append("<ul>"); in_ul = True
            html.append(p)
        else:
            if in_ul:
                html.append("</ul>"); in_ul = False
            html.append(p)
    if in_ul:
        html.append("</ul>")
    return "\n".join(html)


def _is_raw_html(content):
    c = (content or "").lstrip().lower()
    return c.startswith("<!doctype") or c.startswith("<html") or "<body" in c


def _build_full_html(title, content):
    if _is_raw_html(content):
        return content
    body = _lines_to_html_body(title, content)
    return (
        "<!DOCTYPE html>\n<html lang=\"zh\">\n<head>\n<meta charset=\"utf-8\">\n"
        "<title>%s</title>\n"
        "<style>body{font-family:'Noto Sans CJK SC','DejaVu Sans',sans-serif;"
        "margin:40px;line-height:1.6;color:#222}h1{border-bottom:2px solid #444}"
        "h2{color:#333;margin-top:1.2em}</style>\n</head>\n<body>\n%s\n</body>\n</html>"
        % (_esc(title or "Document"), body)
    )


def _create(fmt, base_dir):
    filename = params.get("filename") or ("output.%s" % fmt)
    if not filename.lower().endswith("." + fmt):
        filename = filename + "." + fmt
    out_path = filename if os.path.isabs(filename) else os.path.join(base_dir, filename)
    title = params.get("title")
    content = params.get("content") or ""
    result = {"ok": False, "fmt": fmt}

    if fmt == "docx":
        from docx import Document
        doc = Document()
        if title:
            doc.add_heading(title, 0)
        for raw in content.split("\n"):
            line = raw.rstrip()
            if not line.strip():
                continue
            if line.startswith("## "):
                doc.add_heading(line[3:].strip(), level=2)
            elif line.startswith("# "):
                doc.add_heading(line[2:].strip(), level=1)
            elif line.startswith("- ") or line.startswith("* "):
                doc.add_paragraph(line[2:].strip(), style="List Bullet")
            else:
                doc.add_paragraph(line)
        doc.save(out_path)

    elif fmt == "html":
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(_build_full_html(title, content))

    elif fmt == "pdf":
        html = _build_full_html(title, content)
        html_tmp = out_path[:-4] + ".src.html"
        with open(html_tmp, "w", encoding="utf-8") as f:
            f.write(html)
        r = subprocess.run(["/usr/local/bin/weasyprint", html_tmp, out_path],
                           capture_output=True, text=True)
        if r.returncode != 0 or not os.path.exists(out_path):
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt
            from matplotlib.backends.backend_pdf import PdfPages
            lines = []
            if title:
                lines.append(title); lines.append("")
            lines += content.split("\n")
            with PdfPages(out_path) as pdf:
                per = 42
                for i in range(0, max(1, len(lines)), per):
                    chunk = lines[i:i + per] or [""]
                    fig = plt.figure(figsize=(8.27, 11.69))
                    fig.text(0.08, 0.95, "\n".join(chunk), va="top", ha="left",
                             fontsize=11, family="DejaVu Sans")
                    pdf.savefig(fig); plt.close(fig)
            result["pdf_engine"] = "matplotlib(fallback): " + (r.stderr or "")[:160]
        else:
            result["pdf_engine"] = "weasyprint"

    elif fmt == "pptx":
        from pptx import Presentation
        prs = Presentation()
        slides = params.get("slides")
        if not slides:
            slides = []
            cur = None
            for raw in content.split("\n"):
                line = raw.rstrip()
                if line.startswith("# "):
                    cur = {"title": line[2:].strip(), "bullets": []}
                    slides.append(cur)
                elif line.strip():
                    if cur is None:
                        cur = {"title": title or "Slide", "bullets": []}
                        slides.append(cur)
                    b = line[2:].strip() if (line.startswith("- ") or line.startswith("* ")) else line.strip()
                    cur["bullets"].append(b)
            if not slides:
                slides = [{"title": title or "Slide", "bullets": [content.strip() or "(empty)"]}]
        for s in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = s.get("title", "")
            body = slide.placeholders[1].text_frame
            bullets = s.get("bullets") or []
            body.text = bullets[0] if bullets else ""
            for extra in bullets[1:]:
                body.add_paragraph().text = extra
        prs.save(out_path)
    else:
        result["error"] = "unsupported format: %s" % fmt
        return result

    result["ok"] = os.path.exists(out_path)
    result["path"] = out_path
    result["size_bytes"] = os.path.getsize(out_path) if os.path.exists(out_path) else 0
    return result


def _read(fmt, base_dir):
    raw_path = params.get("path")
    path = raw_path if os.path.isabs(raw_path or "") else os.path.join(base_dir, raw_path or "")
    result = {"ok": False}
    if not os.path.exists(path):
        result["error"] = "file not found: %s" % path
        return result
    if not fmt:
        fmt = path.rsplit(".", 1)[-1].lower()
    text = ""
    meta = {}
    try:
        if fmt == "docx":
            from docx import Document
            d = Document(path)
            text = "\n".join(p.text for p in d.paragraphs if p.text)
            meta["paragraphs"] = len(d.paragraphs)
        elif fmt == "pptx":
            from pptx import Presentation
            prs = Presentation(path)
            chunks = []
            for sl in prs.slides:
                for sh in sl.shapes:
                    if sh.has_text_frame:
                        for para in sh.text_frame.paragraphs:
                            if para.text:
                                chunks.append(para.text)
            text = "\n".join(chunks)
            meta["slides"] = len(list(prs.slides))
        elif fmt == "html":
            from html.parser import HTMLParser
            class _P(HTMLParser):
                def __init__(self):
                    super().__init__(); self.buf = []; self.skip = 0
                def handle_starttag(self, tag, attrs):
                    if tag in ("style", "script"):
                        self.skip += 1
                def handle_endtag(self, tag):
                    if tag in ("style", "script") and self.skip > 0:
                        self.skip -= 1
                def handle_data(self, d):
                    if self.skip == 0 and d.strip():
                        self.buf.append(d.strip())
            p = _P()
            with open(path, encoding="utf-8") as f:
                p.feed(f.read())
            text = "\n".join(p.buf)
        elif fmt == "pdf":
            from pypdf import PdfReader
            rd = PdfReader(path)
            text = "\n".join((pg.extract_text() or "") for pg in rd.pages)
            meta["pages"] = len(rd.pages)
        else:
            result["error"] = "unsupported format: %s" % fmt
            return result
    except Exception as e:
        result["error"] = "read failed: %r" % e
        return result
    result.update({"ok": True, "path": path, "fmt": fmt,
                   "text": text, "text_len": len(text), "meta": meta})
    return result


def _main():
    action = params.get("action")
    fmt = (params.get("fmt") or "").lower()
    base_dir = params.get("base_dir") or "/home/gem/veadk_docs"
    os.makedirs(base_dir, exist_ok=True)
    if action == "create":
        return _create(fmt, base_dir)
    if action == "read":
        return _read(fmt, base_dir)
    return {"ok": False, "error": "unknown action: %s" % action}


print("DOC_RESULT_JSON:" + json.dumps(_main(), ensure_ascii=False))
'''


def _sandbox_session_id(tool_context: Optional[ToolContext]) -> str:
    """与 run_sandbox_agent 一致地基于会话派生 session id，保证同一会话文件可复用。"""
    if tool_context is not None:
        try:
            ic = tool_context._invocation_context
            return f"{ic.agent.name}_{ic.user_id}_{ic.session.id}"
        except Exception:  # noqa: BLE001
            pass
    return os.getenv("VEADK_DOC_TOOLS_SESSION", "veadk-doc-tools")


def _extract_sandbox_result(res: dict) -> dict:
    """从 InvokeTool/RunCode 返回体里抽取脚本 stdout，并解析成 JSON dict。"""
    # 先识别 OpenAPI 错误信封（如会话数超限 CreateSessionFailed），给出可读错误。
    if isinstance(res, dict) and "Result" not in res:
        api_error = (res.get("ResponseMetadata") or {}).get("Error") or {}
        if api_error:
            return {
                "ok": False,
                "error": f"AgentKit {api_error.get('Code', 'Error')}: "
                f"{api_error.get('Message', '')}".strip(),
            }
        return {"ok": False, "error": f"unexpected sandbox response: {str(res)[:300]}"}

    try:
        payload = json.loads(res["Result"]["Result"])
    except Exception as e:  # noqa: BLE001
        return {"ok": False, "error": f"invalid sandbox response: {e!r}"}

    stdout_text = ""
    exc = ""
    for out in payload.get("data", {}).get("outputs", []) or []:
        if not isinstance(out, dict):
            continue
        if out.get("text"):
            stdout_text += out["text"]
        if out.get("ename"):
            exc = f"{out.get('ename')}: {out.get('evalue')}"

    # 优先解析脚本输出的结果标记行（对 kernel 的 SystemExit 等噪音免疫）。
    marker = "DOC_RESULT_JSON:"
    for line in stdout_text.splitlines():
        idx = line.find(marker)
        if idx != -1:
            try:
                return json.loads(line[idx + len(marker):].strip())
            except Exception:  # noqa: BLE001
                continue

    if not payload.get("success", True):
        return {"ok": False, "error": payload.get("message", "sandbox run failed"),
                "stdout": stdout_text[:400]}
    if exc:
        return {"ok": False, "error": f"sandbox exception: {exc}", "stdout": stdout_text}

    return {"ok": False, "error": "no parseable JSON from sandbox", "stdout": stdout_text[:400]}


def _run_doc_script(params: dict, tool_context: Optional[ToolContext]) -> dict:
    params.setdefault("base_dir", _SANDBOX_DOC_DIR)
    params_b64 = base64.b64encode(
        json.dumps(params, ensure_ascii=False).encode("utf-8")
    ).decode("ascii")
    code = _SANDBOX_SCRIPT.replace("__PARAMS_B64__", params_b64)

    tool_id = resolve_agentkit_tool_id("AGENTKIT_TOOL_ID_SCRIPT")
    session_id = _sandbox_session_id(tool_context)
    tool_state = tool_context.state if tool_context is not None else None

    logger.info(
        "doc-tool sandbox run: action=%s fmt=%s tool_id=%s session=%s",
        params.get("action"), params.get("fmt"), tool_id, session_id,
    )
    res = invoke_agentkit_run_code(
        tool_id=tool_id,
        tool_user_session_id=session_id,
        code=code,
        timeout=_RUN_TIMEOUT,
        kernel_name="python3",
        tool_state=tool_state,
    )
    return _extract_sandbox_result(res)


def create_document(
    doc_format: str,
    filename: str,
    content: str,
    title: Optional[str] = None,
    tool_context: Optional[ToolContext] = None,
) -> dict:
    """Generate a Word / PDF / PPT / HTML document inside the AgentKit sandbox.

    The file is created with the image's built-in libraries (python-docx,
    python-pptx, pypdf, weasyprint) and stored under /home/gem/veadk_docs so it
    can later be read back with read_document.

    Args:
        doc_format: One of "docx", "pdf", "pptx", "html".
        filename: Target file name (extension optional, will be normalized).
        content: The document body text. Lines beginning with "# "/"## " become
            headings/slide titles; lines beginning with "- "/"* " become bullets.
            For HTML you may also pass a full raw HTML string.
        title: Optional document/first-heading title.
    Returns:
        A dict: {ok, path, fmt, size_bytes, ...} on success, or {ok: False, error}.
    """
    fmt = (doc_format or "").lower().strip()
    if fmt not in _SUPPORTED_FORMATS:
        return {"ok": False, "error": f"unsupported doc_format {doc_format!r}; "
                f"choose one of {_SUPPORTED_FORMATS}"}
    if not (content or "").strip() and not title:
        return {"ok": False, "error": "content and title are both empty"}

    params = {
        "action": "create",
        "fmt": fmt,
        "filename": filename,
        "content": content or "",
        "title": title,
    }
    result = _run_doc_script(params, tool_context)
    return result


def read_document(
    path: str,
    doc_format: Optional[str] = None,
    tool_context: Optional[ToolContext] = None,
) -> dict:
    """Read back a Word / PDF / PPT / HTML file from the AgentKit sandbox.

    Extracts plain text so the agent can verify or summarize the file content.

    Args:
        path: File path or bare file name previously created under
            /home/gem/veadk_docs (bare names are resolved against that dir).
        doc_format: Optional explicit format ("docx"/"pdf"/"pptx"/"html");
            inferred from the extension when omitted.
    Returns:
        A dict: {ok, path, fmt, text, text_len, meta} on success, else {ok: False, error}.
    """
    if not (path or "").strip():
        return {"ok": False, "error": "path is empty"}
    fmt = (doc_format or "").lower().strip()
    if fmt and fmt not in _SUPPORTED_FORMATS:
        return {"ok": False, "error": f"unsupported doc_format {doc_format!r}"}

    params = {
        "action": "read",
        "fmt": fmt,
        "path": path,
    }
    return _run_doc_script(params, tool_context)
